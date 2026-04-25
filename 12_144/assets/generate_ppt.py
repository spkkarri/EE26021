"""
Generate a professional PowerPoint presentation from the self-driving car project outputs.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image
import glob
from datetime import datetime

# ----------------------------------------------------------------------
# Helper functions to create slides
# ----------------------------------------------------------------------
def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle

def add_bullet_slide(prs, title, bullets, notes=""):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = bullets[0]
    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

def add_image_slide(prs, title, image_path, caption=""):
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    img = Image.open(image_path)
    aspect = img.height / img.width
    height = width * aspect
    slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.size = Pt(24)
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(1), top + height + Inches(0.2), Inches(8), Inches(0.5))
        cap_box.text_frame.text = caption
        cap_box.text_frame.paragraphs[0].font.size = Pt(12)
        cap_box.text_frame.paragraphs[0].font.italic = True

def add_two_images_slide(prs, title, img1_path, img2_path, caption1="", caption2=""):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.size = Pt(24)
    # Left image
    left1 = Inches(0.5)
    top = Inches(1.5)
    width1 = Inches(4)
    img1 = Image.open(img1_path)
    aspect1 = img1.height / img1.width
    height1 = width1 * aspect1
    slide.shapes.add_picture(img1_path, left1, top, width=width1, height=height1)
    # Right image
    left2 = Inches(5.5)
    width2 = Inches(4)
    img2 = Image.open(img2_path)
    aspect2 = img2.height / img2.width
    height2 = width2 * aspect2
    slide.shapes.add_picture(img2_path, left2, top, width=width2, height=height2)
    # Captions
    if caption1:
        cap1 = slide.shapes.add_textbox(left1, top + height1 + Inches(0.1), width1, Inches(0.5))
        cap1.text_frame.text = caption1
        cap1.text_frame.paragraphs[0].font.size = Pt(10)
    if caption2:
        cap2 = slide.shapes.add_textbox(left2, top + height2 + Inches(0.1), width2, Inches(0.5))
        cap2.text_frame.text = caption2
        cap2.text_frame.paragraphs[0].font.size = Pt(10)

# ----------------------------------------------------------------------
# Main presentation generation
# ----------------------------------------------------------------------
print("Generating PowerPoint presentation...")

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# 1. Title slide
add_title_slide(prs, "Self-Driving Car using Deep Reinforcement Learning", "Project Report – EE2061 Introduction to Machine Learning")

# 2. Team & Guide
add_bullet_slide(prs, "Team & Guide", [
    "Team Members:",
    "  • K Karthik Sai Chaitanya",
    "  • A Sai Teja",
    "  • R Vamsi Meenan",
    "",
    "Guided By: Dr. Phani Krishna",
    "Assistant Professor, Department of Electrical Engineering"
])

# 3. Introduction
add_bullet_slide(prs, "Introduction", [
    "Developed a simulation system for autonomous vehicle movement in a dynamic virtual environment.",
    "Vehicles act as independent agents that observe surroundings and adjust movement.",
    "Combines algorithmic decision logic, mathematical motion modelling, and environment-agent interaction.",
    "Demonstrates realistic traffic simulation using computational techniques."
])

# 4. Problem Statement
add_bullet_slide(prs, "Problem Statement", [
    "Real-world traffic is highly dynamic and unpredictable.",
    "Traditional fixed-rule simulations fail to capture variability and adaptability.",
    "Goal: Design a flexible simulation where vehicles respond dynamically to environment.",
    "Integrate motion calculations, adaptive decision logic, and random variations for natural behaviour."
])

# 5. Objectives
add_bullet_slide(prs, "Objectives", [
    "Design a virtual traffic environment with autonomous agents.",
    "Implement mathematical motion calculations (trigonometry, linear algebra, calculus).",
    "Introduce random variations for realistic, non‑repetitive movement.",
    "Allow intelligent decision‑making and adaptive behaviour.",
    "Create a foundation for future reinforcement learning integration."
])

# 6. System Architecture
add_bullet_slide(prs, "System Architecture", [
    "1. MetaDrive Simulator – high‑fidelity driving environment.",
    "2. Dataset Collector – incremental collection of 500,000 driving samples.",
    "3. PPO Agent – learns optimal driving policy.",
    "4. Domain Randomization – friction, gravity, weather, sensor noise.",
    "5. Curriculum Learning – 5 difficulty stages (straight road → roundabout).",
    "6. Visualisation – GIF, 9 charts, summary report.",
    "7. Live Demo – interactive MetaDrive window."
])

# 7. Dataset Collection
add_bullet_slide(prs, "Dataset Collection", [
    "Collected 500,000 transitions (observation, action, reward, next observation).",
    "Progressive difficulty strategy:",
    "  • 0‑100k: random exploration",
    "  • 100k‑200k: straight driving",
    "  • 200k‑300k: moderate steering",
    "  • 300k‑400k: aggressive steering",
    "  • 400k‑500k: expert varied driving",
    "Checkpoints saved every 50,000 samples – supports resuming.",
    "Final dataset size: ~545 MB."
])

# 8. Training Process
add_bullet_slide(prs, "Training Process", [
    "Algorithm: PPO (Proximal Policy Optimization) – stable, efficient.",
    "Domain Randomization: friction, gravity, mass, road friction, lidar noise, weather, lighting.",
    "Curriculum Learning: 5 stages updated during training.",
    "Training steps: 200,000 (approx. 30‑40 minutes on CPU).",
    "Hyperparameters: learning rate 0.0003, batch size 64, n_steps 2048.",
    "Model saved as .zip file."
])

# 9. Performance GIF (first frame as static image)
gif_files = glob.glob("videos/performance_graph_*.gif") + glob.glob("videos/ppo_performance_*.gif")
if gif_files:
    gif_path = sorted(gif_files)[-1]
    from PIL import Image
    img = Image.open(gif_path)
    first_frame = img.convert('RGB')
    png_path = gif_path.replace('.gif', '_firstframe.png')
    first_frame.save(png_path)
    add_image_slide(prs, "Performance GIF (Animated)", png_path, "The GIF shows reward, steering, throttle, and cumulative reward over time. (Static preview – the actual file is animated.)")
else:
    add_bullet_slide(prs, "Performance GIF", ["No GIF found. Please run train_and_animate.py first."])

# 10. Analysis Charts (9‑panel)
chart_files = glob.glob("videos/analysis_charts_*.png") + glob.glob("videos/ppo_charts_*.png")
if chart_files:
    chart_path = sorted(chart_files)[-1]
    add_image_slide(prs, "Analysis Charts (9 subplots)", chart_path, "Includes rewards, cumulative reward, steering/throttle distributions, action space, loss curves, algorithm comparison, and training details.")
else:
    add_bullet_slide(prs, "Analysis Charts", ["No charts found. Run train_and_animate.py first."])

# 11. Summary Report (excerpt)
summary_files = glob.glob("videos/summary_report_*.txt") + glob.glob("videos/summary_*.txt")
if summary_files:
    with open(sorted(summary_files)[-1], 'r') as f:
        summary_text = f.read()[:800]
    add_bullet_slide(prs, "Summary Report (Excerpt)", summary_text.split('\n')[:15])
else:
    add_bullet_slide(prs, "Summary Report", ["No summary report found."])

# 12. Live Demo
add_bullet_slide(prs, "Live Demo", [
    "Run `python video.py` to watch the trained car drive in real time.",
    "A MetaDrive window opens with rendering enabled.",
    "The car drives until it crashes or reaches destination.",
    "Console prints step‑by‑step reward, total reward, and final statistics.",
    "Demonstrates the trained policy in action."
])

# 13. Results and Observations
add_bullet_slide(prs, "Results and Observations", [
    "Successfully trained a PPO agent that drives on multiple road types.",
    "Average episode length: ~150‑200 steps (depending on training).",
    "Mean reward per step: positive (avoiding crashes, staying on road).",
    "Steering centred near zero, throttle mostly positive – good driving behaviour.",
    "Domain randomization and curriculum improved robustness.",
    "All outputs (GIF, charts, report) generated automatically."
])

# 14. Conclusion
add_bullet_slide(prs, "Conclusion", [
    "Developed a complete self‑driving car pipeline using reinforcement learning.",
    "Collected 500,000 diverse driving samples.",
    "Trained a PPO agent with domain randomization and curriculum.",
    "Generated professional visualisations and a live demo.",
    "Project is research‑grade and can be extended with more advanced algorithms (SAC, Transformer, hierarchical options).",
    "Provides a strong foundation for future work in autonomous driving."
])

# 15. Future Work
add_bullet_slide(prs, "Future Work", [
    "Implement Soft Actor‑Critic (SAC) for higher sample efficiency.",
    "Add transformer‑based temporal memory (sliding window).",
    "Incorporate hierarchical options (meta‑controller + skills).",
    "Deploy on real hardware (Donkey Car / F1TENTH) using domain adaptation.",
    "Integrate with real‑time sensor data from cameras and LiDAR."
])

# 16. Thank You
add_title_slide(prs, "Thank You", "Questions?")

# Save the presentation
output_pptx = "Self_Driving_Car_Presentation.pptx"
prs.save(output_pptx)
print(f"\n✅ Presentation saved as {output_pptx}")
print(f"   Location: {os.path.abspath(output_pptx)}")